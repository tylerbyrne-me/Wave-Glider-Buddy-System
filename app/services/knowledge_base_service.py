"""
Knowledge Base Service

Handles document processing, text extraction, and search logic.
"""

import logging
from pathlib import Path
from typing import Optional
import asyncio

logger = logging.getLogger(__name__)


class KnowledgeBaseService:
    """Service for knowledge base operations."""
    
    def __init__(self):
        """Initialize the knowledge base service."""
        pass
    
    async def extract_text_from_document(
        self,
        file_path: Path,
        file_type: str
    ) -> str:
        """
        Extract searchable text from a document.
        
        Args:
            file_path: Path to the document file
            file_type: Type of file (pdf, docx, pptx)
            
        Returns:
            Extracted text content
        """
        try:
            if file_type == "pdf":
                return await self._extract_pdf_text(file_path)
            elif file_type == "docx":
                return await self._extract_docx_text(file_path)
            elif file_type == "pptx":
                return await self._extract_pptx_text(file_path)
            else:
                logger.warning(f"Unsupported file type for text extraction: {file_type}")
                return ""
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return ""
    
    async def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF."""
        try:
            import PyPDF2
            # Run in thread pool to avoid blocking
            loop = asyncio.get_event_loop()
            text = await loop.run_in_executor(None, self._extract_pdf_sync, file_path)
            return text
        except ImportError:
            logger.warning("PyPDF2 not installed, skipping PDF text extraction")
            return ""
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""
    
    def _extract_pdf_sync(self, file_path: Path) -> str:
        """Synchronous PDF extraction (runs in thread pool)."""
        import PyPDF2
        text = ""
        with open(file_path, "rb") as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page in pdf_reader.pages:
                try:
                    text += page.extract_text() + "\n"
                except Exception as e:
                    logger.warning(f"Error extracting text from PDF page: {e}")
        return text
    
    async def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document
            loop = asyncio.get_event_loop()
            text = await loop.run_in_executor(None, self._extract_docx_sync, file_path)
            return text
        except ImportError:
            logger.warning("python-docx not installed, skipping DOCX text extraction")
            return ""
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""
    
    def _extract_docx_sync(self, file_path: Path) -> str:
        """Synchronous DOCX extraction (runs in thread pool)."""
        from docx import Document
        doc = Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    
    async def _extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation
            loop = asyncio.get_event_loop()
            text = await loop.run_in_executor(None, self._extract_pptx_sync, file_path)
            return text
        except ImportError:
            logger.warning("python-pptx not installed, skipping PPTX text extraction")
            return ""
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    def _extract_pptx_sync(self, file_path: Path) -> str:
        """Synchronous PPTX extraction (runs in thread pool)."""
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

